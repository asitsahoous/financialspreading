"""Extract searchable markdown from all documents in this workspace."""

from __future__ import annotations

import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
INDEX_DIR = ROOT / "_index"


def clean(text: str | None) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def extract_docx(path: Path) -> list[str]:
    with zipfile.ZipFile(path) as z:
        xml = z.read("word/document.xml")
    root = ET.fromstring(xml)
    texts: list[str] = []
    for p in root.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p"):
        parts = [
            t.text
            for t in p.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t")
            if t.text
        ]
        line = clean("".join(parts))
        if line:
            texts.append(line)
    return texts


def extract_pptx(path: Path) -> list[tuple[int, list[str]]]:
    prs = Presentation(str(path))
    slides: list[tuple[int, list[str]]] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(clean(shape.text))
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(clean(cell.text) for cell in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)
        if parts:
            slides.append((i, parts))
    return slides


def extract_pdf(path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(path))
    pages: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages, 1):
        text = clean(page.extract_text() or "")
        if text:
            pages.append((i, text))
    return pages


def write_md(out_path: Path, title: str, source: str, body_lines: list[str]) -> None:
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(f"# {title}\n\n")
        f.write(f"**Source:** `{source}`\n\n")
        f.write("---\n\n")
        f.write("\n\n".join(body_lines))


def main() -> None:
    INDEX_DIR.mkdir(exist_ok=True)
    manifest: list[tuple[str, str, str, str | int]] = []

    for path in sorted(ROOT.rglob("*")):
        if not path.is_file():
            continue
        rel = path.relative_to(ROOT)
        rel_str = rel.as_posix()
        if rel_str.startswith("_index/") or rel_str.startswith("scripts/"):
            continue
        if rel_str.startswith("temp_") or rel_str.endswith(".zip"):
            continue

        ext = path.suffix.lower()
        safe = re.sub(r"[^a-zA-Z0-9._-]+", "_", rel_str)
        out = INDEX_DIR / f"{safe}.md"

        try:
            if ext == ".docx":
                lines = extract_docx(path)
                write_md(out, path.stem, rel_str, lines)
                manifest.append((rel_str, out.name, "docx", len(lines)))
            elif ext == ".pptx":
                slides = extract_pptx(path)
                body = [f"## Slide {n}\n\n" + "\n\n".join(parts) for n, parts in slides]
                write_md(out, path.stem, rel_str, body)
                manifest.append((rel_str, out.name, "pptx", len(slides)))
            elif ext == ".pdf":
                pages = extract_pdf(path)
                body = [f"## Page {n}\n\n{text}" for n, text in pages]
                write_md(out, path.stem, rel_str, body)
                manifest.append((rel_str, out.name, "pdf", len(pages)))
            elif ext == ".md":
                content = path.read_text(encoding="utf-8", errors="ignore")
                out.write_text(
                    f"# {path.stem}\n\n**Source:** `{rel_str}`\n\n---\n\n{content}",
                    encoding="utf-8",
                )
                manifest.append((rel_str, out.name, "md", content.count("\n")))
        except Exception as exc:
            manifest.append((rel_str, "ERROR", ext, str(exc)))

    with open(INDEX_DIR / "MANIFEST.md", "w", encoding="utf-8") as f:
        f.write("# Financial Spreading Board — File Index\n\n")
        f.write("Searchable text exports for all workspace documents.\n\n")
        f.write("| Source File | Indexed As | Type | Units |\n")
        f.write("|-------------|------------|------|-------|\n")
        for src, out, typ, units in manifest:
            f.write(f"| `{src}` | `{out}` | {typ} | {units} |\n")

    print(f"Indexed {len(manifest)} files into {INDEX_DIR}")


if __name__ == "__main__":
    main()
